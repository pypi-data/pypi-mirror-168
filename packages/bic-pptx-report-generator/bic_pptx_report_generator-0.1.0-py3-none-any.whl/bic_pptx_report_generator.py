import pptx
from pptx.util import Pt, Inches
import pyautogui as pa
import os

#make functions
def confirmations():
    #confirm all necessary elements exist
    answer = pa.confirm(text='Is your project folder in the Projects folder of the H drive or the Test Profiles folder?', title='', 
        buttons=['Projects', 'Tests', 'Neither'])
    if answer != 'Neither':
        if answer == 'Projects':
            outside_folder = 'H:\\Projects'
        elif answer == 'Tests':
            outside_folder = 'H:\\Tests\\Test Profiles'
        answer = pa.confirm(text='Does that project have a report folder?', buttons=['Yes', 'No'])
        if answer == 'Yes':
            answer = pa.confirm(text='Does your report folder follow the folder diagram from the H drive?', buttons=['Yes', 'No'])
            if answer == 'Yes':
                answer = pa.confirm(text='Do your folders have at least all the images listed in the text file in the H drive?', buttons=['Yes', 'No'])
                if answer == 'Yes':
                    answer = pa.confirm(text='Are any descriptions included in .txt files with the same name as the picture?', buttons=['Yes', 'No'])
                    if answer == 'Yes':
                        return(True, outside_folder)
                    else:
                        pa.alert(text = 'You shall not pass!')
                else:
                    pa.alert(text = 'Oooh, you gotta be better than that')
            else:
                pa.alert(text = 'That"s rude. Come back when you aren"t wasting my time.')
        else:
            pa.alert(text = 'One does not simply make a report without a report folder')
    else:
        pa.alert(text = 'There"s your problem. Try again when your better at this.')
        return(False)

def setup_path(outside_folder):
    #get folder paths and check they exist
    proceed = False
    while not proceed:
        project_name = pa.prompt(text = 'Name of project folder: ', title = 'Project Folder Name')
        try:
            project_path = outside_folder
            folder_path = os.path.join(project_path, project_name)
            proceed = True
        except Exception as e:
            pa.alert(text = 'Project folder does not exist in H drive. Check spelling and capitalization.')
    proceed = False
    while not proceed:
        try:
            folder_name = pa.prompt(text='Name of report folder: ', title='Report Folder Name')
            proceed = True
            return [folder_path, os.path.join(folder_path, folder_name), folder_name]
        except Exception as e:
            pa.alert('Report folder does not exist. Check spelling and capitalization.')

def title_slide(folder_name, report):
    #create a title slide
    title_slide_layout = report.slide_layouts[0]
    title_slide = report.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    title.text = folder_name

def add_label(slide, pic_label, left, top, width, height):
    #add label to image
    textbox = slide.shapes.add_textbox(left, top+height, width, height)
    textframe = textbox.text_frame
    textframe.text = pic_label
    textframe.paragraphs[0].font.size = Pt(14)

def add_pic(slide, pic_path, pic_label, left, top, width, height, isLabeled = False, withText = False, text_path = ''):
    #add an image to the slide
    if isLabeled == True:
        add_label(slide, pic_label, left, top, width, height)
    if withText == True:
        lines = read_file(text_path)
        add_text_top(slide, lines, left, top, width)
        pic = slide.shapes.add_picture(pic_path, left, top+Inches(1.5), width, height-Inches(2))
    else:
        pic = slide.shapes.add_picture(pic_path, left, top, width, height)

def have_text(directory):
    for str in directory:
        if str.find('.txt') > 0:
            return True

def add_text_right(slide, lines):
    #add text from text file to right side of graph
    count = 0
    textbox = slide.shapes.add_textbox(Inches(11), Inches(1), Inches(2), Inches(5))
    tf = textbox.text_frame
    for line in lines:
        if count == 0 :
            tf.text = line
            count += 1
        else:
            p = tf.add_paragraph()
            p.level = 0
            p.text = line 
            p.font.size = Pt(20)     

def add_text_top(slide, lines, left, top, width):
    #add text from text box
    count = 0
    textbox = slide.shapes.add_textbox(left, top, width, Inches(1.5))
    tf = textbox.text_frame
    for line in lines:
        if count == 0 :
            tf.text = line
            count += 1
        else:
            p = tf.add_paragraph()
            p.level = 0
            p.text = line
    return Inches(1.5)

def read_file(text_path):
    lines = []
    with open(text_path, 'r') as f:
        for line in f:
            lines.append(line)
    f.close()
    return lines

def slide4pic_uneven(slide, pic_num, pic_path, pic_label, isLabeled = False, hasText = False, text_path = ''):
    match(pic_num):
        case 1:
            left = Inches(0.5)
            top = Inches(1)
            width = Inches(2)
            height = Inches(2)
        case 2:
            left = Inches(4)
            top = Inches(1)
            width = Inches(3)
            height = Inches(6)
        case 3:
            left = Inches(8.5)
            top = Inches(2.5)
            width = Inches(4)
            height = Inches(2.5)
        case 4:
            left = Inches(0.5)
            top = Inches(4) 
            width = Inches(2) 
            height = Inches(2) 
    add_pic(slide, pic_path, pic_label, left, top, width, height, isLabeled, hasText, text_path)

def slide4pic_even(slide, pic_num, pic_path, pic_label, isLabeled = False, hasText = False, text_path = ''):
    match(pic_num):
        case 1:
            left = Inches(0.5) 
            top = Inches(1) 
            width = Inches(5) 
            height = Inches(2.5)
        case 2:
            left = Inches(6.5) 
            top = Inches(1) 
            width = Inches(5) 
            height = Inches(2.5) 
        case 3:
            left = Inches(0.5) 
            top = Inches(4) 
            width = Inches(5) 
            height = Inches(2.5) 
        case 4:
            left = Inches(6.5) 
            top = Inches(4) 
            width = Inches(5) 
            height = Inches(2.5) 
    add_pic(slide, pic_path, pic_label, left, top, width, height, isLabeled, hasText, text_path)

def slide3pic_uneven(slide, pic_num, pic_path, pic_label, isLabeled, text_path = ''):
    match(pic_num):
        case 1:
            left = Inches(0.5) 
            top = Inches(1.5) 
            width = Inches(3) 
            height = Inches(2) 
        case 2:
            left = Inches(4) 
            top = Inches(1.5) 
            width = Inches(4.5) 
            height = Inches(5) 
        case 3:
            left = Inches(9) 
            top = Inches(1.5) 
            width = Inches(4) 
            height = Inches(4) 
    add_pic(slide, pic_path, pic_label, left, top, width, height, isLabeled, text_path)

def slide3pic_even(slide, pic_num, pic_path, pic_label, isLabeled = False, hasText = False, text_path = ''):
    match(pic_num):
        case 1:
            left = Inches(0.5) 
            top = Inches(1) 
            width = Inches(3.5) 
            height = Inches(5.2) 
        case 2:
            left = Inches(4.8) 
            top = Inches(1) 
            width = Inches(3.5) 
            height = Inches(5.2) 
        case 3:
            left = Inches(9.1) 
            top = Inches(1) 
            width = Inches(3.5) 
            height = Inches(5.2) 
    add_pic(slide, pic_path, pic_label, left, top, width, height, isLabeled, hasText, text_path)

def slide2pic_even(slide, pic_num, pic_path, pic_label, isLabeled = False, hasText = False, text_path = ''):
    match(pic_num):
        case 1:
            left = Inches(0.5) 
            top = Inches(1) 
            width = Inches(5.5) 
            height = Inches(5.2) 
        case 2:
            left = Inches(7) 
            top = Inches(1) 
            width = Inches(5.5) 
            height = Inches(5.2) 
    add_pic(slide, pic_path, pic_label, left, top, width, height, isLabeled, hasText, text_path)

def slide1_pic_center(slide, pic_path, pic_label, isLabeled = False, hasText = False, text_path = ''):
    left = Inches(3) 
    top = Inches(1) 
    width = Inches(4.5) 
    height = Inches(5) 
    add_pic(slide, pic_path, pic_label, left, top, width, height, isLabeled, hasText, text_path)

def slide_graph(slide, pic_path):
    left = Inches(0.5) 
    top = Inches(1) 
    width = Inches(12) 
    height = Inches(6) 
    add_pic(slide, pic_path, 'None', left, top, width, height)

def slide_graph_words(slide, pic_path, text_path):
    left = Inches(0.5) 
    top = Inches(1) 
    width = Inches(10) 
    height = Inches(6) 
    add_pic(slide, pic_path, 'None', left, top, width, height)
    lines = read_file(text_path)
    add_text_right(slide, lines)

def get_pic_count(directories):
    pic_counter = 0
    for str in directories:
        if str.find('.png') > 0 or str.find('.jpeg') > 0:
            pic_counter += 1
    return pic_counter

def strip_end(str):
    str = str.replace('.jpeg','')
    str = str.replace('.png','')
    str = str.replace('.txt','')
    return str

def make_slide(report):
    slide_layout = report.slide_layouts[5]
    slide = report.slides.add_slide(slide_layout)
    shapes = slide.shapes
    return [slide, shapes]

def text_pic_set(directory):
    pics_with_text = []
    text_files = []
    pics_no_text = []
    for str in directory:
        if str.find('.txt') >= 0:
            for str2 in directory:
                if strip_end(str) == strip_end(str2) and (str2.find('.png') >= 0 or str2.find('.jpeg') >= 0):
                    text_files.append(str)
                    pics_with_text.append(str2)
    for str in directory:
        if str in text_files:
            print()
        elif str in pics_with_text:
            print()
        else:
            pics_no_text.append(str)
    return [pics_with_text, text_files, pics_no_text]

def test_setup_slide(report, path1):
    directories1 = os.listdir(path1)
    has_text = have_text(directories1)
    for str in directories1:
        if str == 'Test Setup':
            [test_setup_slide, shapes] = make_slide(report)
            shapes.title.text = 'Test Setup'
            path2 = os.path.join(path1, 'Test Setup')
            directories2 = os.listdir(path2)
            pic_counter = get_pic_count(directories2)
            if pic_counter > 4:
                pa.alert(text='Too many pictures for test setup. Limit to: DAQ, cell in setup, TC placement, and one other')
            elif pic_counter <=2:
                pa.alert(text='Too few pictures for test setup. Make sure you have: DAQ, cell in setup, and TC placement.')
            if pic_counter == 4:
                for str in directories2:
                    image_path = os.path.join(path2, str)
                    if str.find('DAQ') >= 0:
                        slide4pic_uneven(test_setup_slide, 1, image_path, 'Data Aquisition', True)
                    elif str.find('Cell') >= 0:
                        slide4pic_uneven(test_setup_slide, 2, image_path, 'Cell Setup', True)
                    elif str.find('TC') >= 0:
                        slide4pic_uneven(test_setup_slide, 3, image_path, 'Thermocouple Placement', True)
                    else:
                        str = strip_end(str)
                        slide4pic_uneven(test_setup_slide, 4, image_path, str, True)
            if pic_counter == 3:
                for str in directories2:
                    image_path = os.path.join(path2, str)
                    if str.find('DAQ') >= 0:
                        slide3pic_uneven(test_setup_slide, 1, image_path, 'Data Aquisition Device', True)
                    elif str.find('Cell') >= 0:
                        slide3pic_uneven(test_setup_slide, 2, image_path, 'Cell Setup', True)
                    elif str.find('TC') >= 0:
                        slide3pic_uneven(test_setup_slide, 3, image_path, 'Thermocouple Placement', True)  
    
def cell_slides(report, path1):
    directories1 = os.listdir(path1)
    for str in directories1:
        if str.find('Cell') >= 0:
            path2 = os.path.join(path1, str)
            directories2 = os.listdir(path2)
            for str in directories2:
                if str.find('Graphs') >= 0:
                    path3 = os.path.join(path2, str)
                    directories3 = os.listdir(path3)
                    if have_text(directories3):
                        [graphs_with_text, text_files, graphs_no_text] = text_pic_set(directories3)
                        for index in range(len(text_files)):
                                text_path = os.path.join(path3, text_files[index])
                                graph_path = os.path.join(path3, graphs_with_text[index])
                                [data_slide, shapes] = make_slide(report)
                                slide_graph_words(data_slide, graph_path, text_path)
                                str = strip_end(graphs_with_text[index])
                                shapes = data_slide.shapes
                                shapes.title.text = str
                        for str in graphs_no_text:
                            [data_slide, shapes] = make_slide(report)
                            graph_path = os.path.join(path3, str)
                            slide_graph(data_slide, graph_path)
                            str = strip_end(str)
                            shapes.title.text = strip_end(str)
                    else:
                        for str in directories3:
                            [data_slide, shapes] = make_slide(report)
                            graph_path = os.path.join(path3, str)
                            slide_graph(data_slide, graph_path)
                            str = strip_end(str)
                            shapes.title.text = str
                else:
                    path3 = os.path.join(path2, str)
                    directories3 = os.listdir(path3)
                    [cell_slide, shapes] = make_slide(report)
                    shapes.title.text = strip_end(str)
                    pic_counter = get_pic_count(directories3)
                    if pic_counter > 4:
                        print('Too many before pictures. Please have max of 4')
                    elif pic_counter < 1:
                        print('Please add before picture.')
                    elif pic_counter == 4:
                        count = 0
                        if have_text(directories3):
                            [pics_with_text, text_files, pics_no_text] = text_pic_set(directories3)
                            for str in pics_no_text:
                                image_path = os.path.join(path3, str)
                                str = strip_end(str)
                                count += 1
                                slide4pic_even(cell_slide, count, image_path, str)
                            for index in range(len(text_files)):
                                text_path = os.path.join(path3, text_files[index])
                                pic_path = os.path.join(path3, pics_with_text[index])
                                str = strip_end(text_files[index])
                                slide4pic_even(cell_slide, count, pic_path, str, False, True, text_path)
                        else:
                            for str in directories3:
                                image_path = os.path.join(path3, str)
                                str = strip_end(str)
                                count += 1
                                slide4pic_even(cell_slide, count, image_path, str)
                    elif pic_counter == 3:
                        count = 0
                        if have_text(directories3):
                            [pics_with_text, text_files, pics_no_text] = text_pic_set(directories3)
                            for str in pics_no_text:
                                image_path = os.path.join(path3, str)
                                str = strip_end(str)
                                count += 1
                                slide3pic_even(cell_slide, count, image_path, str)
                            for index in range(len(text_files)):
                                text_path = os.path.join(path3, text_files[index])
                                pic_path = os.path.join(path3, pics_with_text[index])
                                str = strip_end(text_files[index])
                                slide3pic_even(cell_slide, count, pic_path, str, False, True, text_path)
                        else:
                            for str in directories3:
                                image_path = os.path.join(path3, str)
                                str = strip_end(str)
                                count += 1
                                slide3pic_even(cell_slide, count, image_path, str)
                    elif pic_counter == 2:
                        count = 0
                        if have_text(directories3):
                            [pics_with_text, text_files, pics_no_text] = text_pic_set(directories3)
                            for str in pics_no_text:
                                image_path = os.path.join(path3, str)
                                str = strip_end(str)
                                count += 1
                                slide2pic_even(cell_slide, count, image_path, str)
                            for index in range(len(text_files)):
                                text_path = os.path.join(path3, text_files[index])
                                pic_path = os.path.join(path3, pics_with_text[index])
                                str = strip_end(text_files[index])
                                slide2pic_even(cell_slide, count, pic_path, str, False, True, text_path)
                        else:
                            for str in directories3:
                                image_path = os.path.join(path3, str)
                                str = strip_end(str)
                                count += 1
                                slide2pic_even(cell_slide, count, image_path, str)
                    elif pic_counter == 1:
                        count = 0
                        if have_text(directories3):
                            [pics_with_text, text_files, pics_no_text] = text_pic_set(directories3)
                            for str in pics_no_text:
                                image_path = os.path.join(path3, str)
                                str = strip_end(str)
                                count += 1
                                slide1_pic_center(cell_slide, count, image_path, str)
                            for index in range(len(text_files)):
                                text_path = os.path.join(path3, text_files[index])
                                pic_path = os.path.join(path3, pics_with_text[index])
                                str = strip_end(text_files[index])
                                slide1_pic_center(cell_slide, pic_path, str, False, True, text_path)
                        else:
                            for str in directories3:
                                image_path = os.path.join(path3, str)
                                str = strip_end(str)
                                count += 1
                                slide1_pic_center(cell_slide, image_path, str)
                


#check all the things are in order
[ready, starting_folder] = confirmations()
if ready:
    #folder path
    [save_path, report_path, folder_name] = setup_path(starting_folder)
    #create a powerpoint
    report = pptx.Presentation(r'C:\Users\Heather.Marsico\Documents\Reporting Standardization\Blank.pptx')
    #make title slide
    title_slide(folder_name, report)
    #make test setup slide
    test_setup_slide(report, report_path)
    #make slides from cell folders
    cell_slides(report, report_path)
    #save
    save_file = os.path.join(save_path, 'Report.pptx')
    report.save(save_file)
    pa.alert(text = 'Yay! I did it! A report has been created in your project folder. Please review it for errors,\
    order of slides, and to make sure it all looks good. If there were any issues, please talk to Spotts.')